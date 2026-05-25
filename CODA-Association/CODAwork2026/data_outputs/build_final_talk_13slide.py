"""
CoDaWork 2026 — Final Talk Deck builder, 13-SLIDE VERSION.

Built 2026-05-24 — expansion from the 10-slide compressed deck.
Per Peter's directive: the per-country navigation chart was too small to read
on the previous slides 6/7/8 (3 charts crammed onto the right margin at
2.6" wide). The 13-slide version splits each country into two slides:
one for the share-and-work view (the existing fig2/3/4), one for the
navigation chart at full size (5.5"+ wide, finally legible from the back).

DOES NOT supersede the 10-slide deck unless promoted; ships as another sibling.

Story arc (13 slides):
   1. Title + question + FULL CONTACT
   2. Size view hides work        (USA Solar 760× hook)
   3. Five viewpoints              (one schematic)
   4. Activation Coefficient       (yeast factor + formula)
   5. Three archetypes overview    (DEU / JPN / GBR)
   6. Germany — share-and-work view       (continuous arc; fig2)
   7. Germany — navigation chart           (course on simplex; fig6_nav_deu LARGE)
   8. Japan — share-and-work view          (shock + reorganisation; fig3)
   9. Japan — navigation chart             (loop on simplex; fig6_nav_jpn LARGE)
  10. UK — share-and-work view            (regime change; fig4)
  11. UK — navigation chart                (jump-and-return on simplex; fig6_nav_gbr LARGE)
  12. 5-of-9 cross-country signature       (fig5)
  13. What the stack answers              (synthesis + AI Use footer)

Under the conference's 15-slide recommendation. Live time target:
~10 min talk + 1 min scroll + 1 min projector + ~3 min Q&A = ~15 min.

Speaking script content unchanged from 10-slide version; the per-country
talk is allocated across the new pair (share-and-work facts on the first
slide of each pair, helmsman-trajectory and course-directness interpretation
on the navigation slide).

MC-4 falsifiability frame remains in manuscript only.
Cinema scroll / projector run after slide 13 verbally; no dedicated slide.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUT = Path("/sessions/epic-gracious-lovelace/mnt/Claude CoWorker/Current-Repo/Hs/CODA-Association/CODAwork2026/data_outputs/CodaWork2026_FinalTalk_13Slide_2026-05-24.pptx")
MANU_FIG = Path("/sessions/epic-gracious-lovelace/mnt/Claude CoWorker/Current-Repo/Hs/papers/codawork2026/manuscript/figures")

# Palette — preserved from 10/22-slide decks for visual continuity
NAVY   = RGBColor(0x0B, 0x1F, 0x33)
GOLD   = RGBColor(0xF2, 0xB6, 0x32)
INK    = RGBColor(0xEE, 0xEE, 0xEE)
DIM    = RGBColor(0xB8, 0xB8, 0xB8)
ACCENT = RGBColor(0xC9, 0x8A, 0x1C)

prs = Presentation()
prs.slide_width  = Inches(11)
prs.slide_height = Inches(8.5)
BLANK = prs.slide_layouts[6]
TOTAL = 13

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    return slide

def add_text(slide, text, left, top, width, height,
             font_size=12, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb

def add_footer(slide, n, label=""):
    add_text(slide, label, 0.5, 8.10, 7.0, 0.30, font_size=9, color=DIM, italic=True)
    add_text(slide, f"{n} / {TOTAL}", 9.5, 8.10, 1.0, 0.30, font_size=9, color=DIM, align=PP_ALIGN.RIGHT)

def add_title_strip(slide, title, subtitle=None):
    add_text(slide, title, 0.5, 0.35, 10.0, 0.70, font_size=26, bold=True, color=INK, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, 0.5, 1.10, 10.0, 0.40, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)

def add_image(slide, path, left, top, width, height=None):
    if height is None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))

# ───────── SLIDE 1 — Title + question + FULL CONTACT ─────────
s = add_slide()
add_text(s, "Compositional monitoring of energy-mix drift on the simplex",
         0.5, 1.0, 10.0, 1.2, font_size=28, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "Which carrier did the structural work?",
         0.5, 2.5, 10.0, 0.6, font_size=18, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Not only which carrier got bigger — which carrier moved the composition.",
         0.5, 3.2, 10.0, 0.5, font_size=13, color=INK, italic=True, align=PP_ALIGN.CENTER)

# Operationalization tagline
add_text(s,
         "Operationalizing compositional analysis — a runnable standard for researchers and the AI assistants they choose.",
         0.5, 4.0, 10.0, 0.4, font_size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# Conference line
add_text(s, "CoDaWork 2026  ·  Coimbra, Portugal  ·  1–5 June 2026",
         0.5, 4.8, 10.0, 0.4, font_size=13, color=DIM, align=PP_ALIGN.CENTER)

# Author + lab
add_text(s, "P. Higgins  ·  Rogue Wave Audio / Binaural Test Lab  ·  Markham, Ontario, Canada",
         0.5, 5.4, 10.0, 0.4, font_size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)

# Contact block — boxed for visual weight
add_text(s, "Contact",
         0.5, 6.05, 10.0, 0.3, font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "PeterHiggins@RogueWaveAudio.com",
         0.5, 6.35, 10.0, 0.35, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "github.com / PeterHiggins19 / higgins-decomposition",
         0.5, 6.75, 10.0, 0.35, font_size=12, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "Community folder:  CODA-Association/   ·   Hand-out in UN-6 locales (EN · FR · ES · RU · ZH · AR)",
         0.5, 7.15, 10.0, 0.35, font_size=11, color=DIM, align=PP_ALIGN.CENTER)

# Doctrine
add_text(s,
         "The instrument reads.   The expert decides.   The hashes carry the receipts.   The vocabulary holds the line.",
         0.5, 7.85, 10.0, 0.35, font_size=10, italic=True, color=DIM, align=PP_ALIGN.CENTER)

# ───────── SLIDE 2 — Size view hides work ─────────
s = add_slide()
add_title_strip(s, "The size view hides the work",
                "A carrier can be small in share and large in structural work.")
add_text(s, "World electricity, 25 years — the standard stacked-area view",
         0.7, 1.85, 9.6, 0.4, font_size=14, bold=True, color=INK)
add_text(s,
         "Coal stays dominant.  Gas grows.  Nuclear declines.\n"
         "Solar appears as a thin yellow sliver after 2010.\n"
         "Wind grows steadily but stays visually small until late.",
         0.7, 2.35, 5.4, 1.6, font_size=13, color=INK)
add_text(s, "What the size view misses",
         6.6, 1.95, 4.0, 0.45, font_size=15, bold=True, color=GOLD)
add_text(s,
         "USA Solar, 2012 → 2013",
         6.6, 2.55, 4.2, 0.45, font_size=14, bold=True, color=INK)
add_text(s,
         "starting share        0.107 %\n"
         "structural Power Share  81.7 %\n"
         "Activation Coefficient   ≈ 760 ×",
         6.6, 3.05, 4.2, 1.5, font_size=14, color=INK, font_name="Consolas")
add_text(s,
         "Solar acted at 760 × its size.\nNo size view shows that.",
         6.6, 4.85, 4.2, 1.0, font_size=13, italic=True, color=ACCENT)
add_text(s,
         "This talk is the reason that number exists.",
         0.7, 6.6, 9.6, 0.5, font_size=14, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
add_text(s,
         "Mathematics: standard compositional data analysis.  Application: monitoring frame may be new.",
         0.7, 7.5, 9.6, 0.4, font_size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 2)

# ───────── SLIDE 3 — Five viewpoints, one observable stack ─────────
s = add_slide()
add_title_strip(s, "Five viewpoints, one observable stack",
                "Each viewpoint answers one question.  Together: an auditable transition event.")
add_image(s, MANU_FIG / "fig1_method.png", 0.6, 1.7, 5.6)
add_text(s, "Composition",      6.5, 1.85, 4.0, 0.40, font_size=15, bold=True, color=GOLD)
add_text(s, "what share each carrier has.", 6.5, 2.25, 4.0, 0.35, font_size=12, color=INK)
add_text(s, "Helmsman",         6.5, 2.75, 4.0, 0.40, font_size=15, bold=True, color=GOLD)
add_text(s, "which carrier has the largest CLR move at a step.", 6.5, 3.15, 4.0, 0.4, font_size=12, color=INK)
add_text(s, "Helmsman trajectory", 6.5, 3.75, 4.0, 0.40, font_size=15, bold=True, color=GOLD)
add_text(s, "when the steering carrier changes.", 6.5, 4.15, 4.0, 0.35, font_size=12, color=INK)
add_text(s, "Power Share",      6.5, 4.65, 4.0, 0.40, font_size=15, bold=True, color=GOLD)
add_text(s, "how much squared CLR motion each carrier did.", 6.5, 5.05, 4.0, 0.4, font_size=12, color=INK)
add_text(s, "Activation Coefficient", 6.5, 5.65, 4.0, 0.40, font_size=15, bold=True, color=GOLD)
add_text(s, "Power Share ÷ starting share — the yeast factor.", 6.5, 6.05, 4.0, 0.4, font_size=12, color=INK)
add_text(s,
         "All five derive from CLR + ILR-Helmert.  Pure CoDa geometry; no new mathematics.",
         0.5, 7.5, 10.0, 0.4, font_size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 3)

# ───────── SLIDE 4 — Activation Coefficient, the yeast factor ─────────
s = add_slide()
add_title_strip(s, "The Activation Coefficient — the yeast factor",
                "How much structural work a carrier does, relative to how much of the mix it is.")
add_text(s, "α_i(t)  =  Power Share_i(t)  ÷  starting share_i(t)",
         0.5, 1.95, 10.0, 0.7, font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font_name="Consolas")
add_text(s, "α ≈ 1     carrier does work proportional to its size — ordinary",
         1.0, 3.0, 9.0, 0.45, font_size=14, color=INK)
add_text(s, "α ≫ 1     carrier acts far above its size — hidden driver",
         1.0, 3.5, 9.0, 0.45, font_size=14, color=GOLD, bold=True)
add_text(s, "α < 1     carrier carries less work than its size suggests — coasting",
         1.0, 4.0, 9.0, 0.45, font_size=14, color=DIM)
add_text(s, "Worked example — USA Solar 2012 → 2013",
         0.7, 4.85, 9.6, 0.45, font_size=15, bold=True, color=INK)
add_text(s,
         "starting share      0.107 %     small\n"
         "Power Share         81.7  %     most of the work\n"
         "α                   ≈ 760 ×     yeast",
         0.7, 5.35, 9.6, 1.5, font_size=14, color=INK, font_name="Consolas")
add_text(s,
         "Yeast is 2% of a loaf by mass and does 100% of the rising. Same shape.",
         0.7, 7.05, 9.6, 0.45, font_size=12, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s,
         "Solar 2010–2015 appears repeatedly as small-share / large-structural-work across the corpus.",
         0.7, 7.55, 9.6, 0.4, font_size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 4)

# ───────── SLIDE 5 — Three archetypes ─────────
s = add_slide()
add_title_strip(s, "Three archetypes — one instrument, three regimes",
                "Same protocol applied to three transitions that look fundamentally different.")
add_text(s, "Germany",   1.05, 1.85, 3.0, 0.5, font_size=18, bold=True, color=GOLD)
add_text(s, "deliberate transition",  1.05, 2.30, 3.0, 0.4, font_size=12, italic=True, color=INK)
add_text(s, "continuous arc",         1.05, 2.65, 3.0, 0.4, font_size=12, italic=True, color=DIM)
add_text(s, "Energiewende\n2000 → 2025\nsolar + wind absorb\nstructural work\nbefore size dominates.",
         1.05, 3.25, 3.0, 3.0, font_size=13, color=INK)

add_text(s, "Japan",     4.5, 1.85, 3.0, 0.5, font_size=18, bold=True, color=GOLD)
add_text(s, "external shock",         4.5, 2.30, 3.0, 0.4, font_size=12, italic=True, color=INK)
add_text(s, "loop and reorganise",    4.5, 2.65, 3.0, 0.4, font_size=12, italic=True, color=DIM)
add_text(s, "Fukushima 2011\ndisplaces nuclear,\ncauses 2011–2013\nmulti-year compositional\nreorganisation.",
         4.5, 3.25, 3.0, 3.0, font_size=13, color=INK)

add_text(s, "United Kingdom", 7.95, 1.85, 3.0, 0.5, font_size=18, bold=True, color=GOLD)
add_text(s, "regime change",          7.95, 2.30, 3.0, 0.4, font_size=12, italic=True, color=INK)
add_text(s, "jump and return",        7.95, 2.65, 3.0, 0.4, font_size=12, italic=True, color=DIM)
add_text(s, "Coal exit 2012–2020\nfrom > 30 % to < 2 %.\nWind, solar, others\nabsorb displaced\nstructural work.",
         7.95, 3.25, 3.0, 3.0, font_size=13, color=INK)

add_text(s,
         "Three different transition regimes.  One operational protocol reads them all.",
         0.5, 7.4, 10.0, 0.5, font_size=13, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
add_footer(s, 5)

# ───────── SLIDE 6 — Germany share-and-work view ─────────
# Layout: case-study figure shrunk to 4.85" tall (was 5.2") to clear bottom for callout + italic + footer
s = add_slide()
add_title_strip(s, "Germany — share-and-work view",
                "Energiewende read as a single smooth arc on the simplex (chart pair: shares + structural work).")
add_image(s, MANU_FIG / "fig2_germany.png", 1.0, 1.55, 9.0, height=4.85)
add_text(s,
         "Solar 2005–2006:  0.21 % share  ·  71.1 % structural work  ·  α ≈ 333 ×",
         0.5, 6.55, 10.0, 0.4, font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Tiny share, dominant work.  The yeast-factor signature, four years before the share view shows anything.",
         0.5, 7.00, 10.0, 0.4, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 6, "Germany — share + work  ·  navigation chart on next slide")

# ───────── SLIDE 7 — Germany navigation chart (LARGE) ─────────
# Layout: nav chart explicit 6.5" × 5.0" to clear bottom; centered horizontally
s = add_slide()
add_title_strip(s, "Germany — course on the simplex",
                "The Helmsman trajectory: where the composition went, year by year.")
add_image(s, MANU_FIG / "fig6_nav_deu.png", 2.25, 1.50, 6.5, height=5.0)
add_text(s,
         "Course directness 0.41  —  continuous arc toward the renewable vertex.",
         0.5, 6.65, 10.0, 0.45, font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Smooth, monotone reorientation — no loops, no flips.  Deliberate transition as a sustained course.",
         0.5, 7.15, 10.0, 0.4, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 7, "Germany — navigation chart")

# ───────── SLIDE 8 — Japan share-and-work view ─────────
s = add_slide()
add_title_strip(s, "Japan — share-and-work view",
                "Fukushima 2011: external shock displaces nuclear, multi-year compositional reorganisation.")
add_image(s, MANU_FIG / "fig3_japan.png", 1.0, 1.55, 9.0, height=4.85)
add_text(s,
         "Aitchison distance 2011 → 2012  ≈ 3 × neighbouring-year baseline  ·  helmsman flips 17 ×",
         0.5, 6.55, 10.0, 0.4, font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "The instrument detects both the shock and the multi-year reorganisation that followed.",
         0.5, 7.00, 10.0, 0.4, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 8, "Japan — share + work  ·  navigation chart on next slide")

# ───────── SLIDE 9 — Japan navigation chart (LARGE) ─────────
s = add_slide()
add_title_strip(s, "Japan — course on the simplex",
                "The Helmsman trajectory: looping reorganisation, not a single step.")
add_image(s, MANU_FIG / "fig6_nav_jpn.png", 2.25, 1.50, 6.5, height=5.0)
add_text(s,
         "Course directness 0.09  —  loop-and-reorganise archetype.",
         0.5, 6.65, 10.0, 0.45, font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Trajectory revisits and reroutes — the system searches for a new composition, not a planned course.",
         0.5, 7.15, 10.0, 0.4, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 9, "Japan — navigation chart")

# ───────── SLIDE 10 — UK share-and-work view ─────────
s = add_slide()
add_title_strip(s, "United Kingdom — share-and-work view",
                "Coal exit as policy-driven regime change, absorbed across multiple renewable carriers.")
add_image(s, MANU_FIG / "fig4_uk.png", 1.0, 1.55, 9.0, height=4.85)
add_text(s,
         "Coal:  > 30 %  →  < 2 %.   Wind, solar, and other renewables absorb the displaced structural work.",
         0.5, 6.55, 10.0, 0.4, font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Displaced work redistributes across several carriers — not one substitute, but several.",
         0.5, 7.00, 10.0, 0.4, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 10, "United Kingdom — share + work  ·  navigation chart on next slide")

# ───────── SLIDE 11 — UK navigation chart (LARGE) ─────────
s = add_slide()
add_title_strip(s, "United Kingdom — course on the simplex",
                "The Helmsman trajectory: a jump from the coal vertex, then a return toward stability.")
add_image(s, MANU_FIG / "fig6_nav_gbr.png", 2.25, 1.50, 6.5, height=5.0)
add_text(s,
         "Course directness 0.36  —  jump-and-return archetype.",
         0.5, 6.65, 10.0, 0.45, font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Sharp departure from the coal vertex, then re-stabilisation — regime change as one decisive displacement.",
         0.5, 7.15, 10.0, 0.4, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 11, "United Kingdom — navigation chart")

# ───────── SLIDE 12 — Cross-country: 5 of 9 ─────────
s = add_slide()
add_title_strip(s, "Cross-country signature — 5 of 9 reproduce the pattern",
                "From three case archetypes to a corpus-level result.")
add_image(s, MANU_FIG / "fig5_crosscountry.png", 1.7, 1.55, 7.6, height=5.0)
add_text(s, "fires   AUS · CHN · GBR · IND · JPN",
         0.5, 7.10, 10.0, 0.45, font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "quiet   DEU (annual) · FRA · USA · WLD",
         0.5, 7.50, 10.0, 0.4, font_size=12, color=DIM, align=PP_ALIGN.CENTER)
add_text(s,
         "A useful detector should not fire everywhere.  Discrimination is itself evidence.",
         0.5, 7.95, 10.0, 0.4, font_size=11, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_footer(s, 12)

# ───────── SLIDE 13 — What the stack answers ─────────
s = add_slide()
add_title_strip(s, "What the stack answers",
                "One observable, five distinct questions, one reproducible object.")
rows = [
    ("WHAT", "carriers are big.",                    "size view"),
    ("WHO",  "is at the wheel.",                     "Helmsman"),
    ("WHEN", "the steering changes.",                "Helmsman trajectory"),
    ("HOW MUCH", "work each carrier did.",           "Power Share"),
    ("WHY",  "a small carrier mattered.",            "Activation Coefficient"),
]
y0 = 2.0
for i, (q, ans, who) in enumerate(rows):
    y = y0 + i * 0.85
    add_text(s, q,   0.9, y, 1.7, 0.6, font_size=20, bold=True, color=GOLD)
    add_text(s, ans, 2.7, y, 4.6, 0.6, font_size=15, color=INK)
    add_text(s, who, 7.3, y, 3.2, 0.6, font_size=13, italic=True, color=DIM)
add_text(s,
         "The stack does not replace interpretation.  It gives interpretation a reproducible object.",
         0.5, 7.5, 10.0, 0.5, font_size=13, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
# AI Use Declaration footer (HUF-STD-001 v1.1 compliance)
add_text(s,
         "AI Use Declaration (HUF-STD-001 v1.1):  research design, mathematical content, code, and scientific "
         "responsibility remain with the named author.  AI assistants (Claude, ChatGPT, Grok) used for drafting, "
         "sweeps, and reviews.  Author retains full responsibility.   Apache-2.0 code  ·  CC BY 4.0 docs.",
         0.5, 7.85, 10.0, 0.6, font_size=8, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 13)
# Save
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
