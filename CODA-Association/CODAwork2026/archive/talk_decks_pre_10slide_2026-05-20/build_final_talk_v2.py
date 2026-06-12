"""
CoDaWork 2026 — Final Talk Deck builder (v2 — slide 12 split into 3 per-country)
Built 2026-05-18, supersedes build_final_talk.py for the 22-slide variant.
Aligned with the manuscript Compositional_Monitoring_2026 and the community study.

Story arc (20 main slides + AI Use Declaration + Standard Stamp = 22):
  1. Title
  2. The Question
  3. The size view hides the driver
  4. Five viewpoints — the operational protocol
  5. Viewpoint 2 — Helmsman
  6. Viewpoint 3 — Power Share
  7. Viewpoint 4 — Activation Coefficient
  8. The hidden driver — solar yeast era headline table
  9. Germany — Energiewende deep-dive
  10. Japan — Fukushima cascade
  11. United Kingdom — coal-exit regime change
  12. Germany — navigation chart (Plate 16, full page)
  13. Japan — navigation chart (Plate 16, full page)
  14. United Kingdom — navigation chart (Plate 16, full page)
  15. Cross-country signature
  16. Synthesis — WHAT + WHY
  17. MC-4 falsifiable claim + four defeat paths
  18. Bridge to data scroll
  19. Bridge to projector + Q&A
  20. Repositories + closing
  21. AI Use Declaration per HUF-STD-001 v1.1
  22. Standard Stamp colophon
"""
import sys
sys.path.insert(0, "/sessions/epic-gracious-lovelace/mnt/Claude CoWorker/Studies/_shared")

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUT = Path("/sessions/epic-gracious-lovelace/mnt/Claude CoWorker/Current-Repo/Hs/CODA-Association/CODAwork2026/data_outputs/CodaWork2026_FinalTalk_2026-05-17.pptx")
MANU_FIG = Path("/sessions/epic-gracious-lovelace/mnt/Claude CoWorker/Current-Repo/Hs/papers/codawork2026/manuscript/figures")

NAVY   = RGBColor(0x0B, 0x1F, 0x33)
GOLD   = RGBColor(0xF2, 0xB6, 0x32)
INK    = RGBColor(0xEE, 0xEE, 0xEE)
DIM    = RGBColor(0xB8, 0xB8, 0xB8)
ACCENT = RGBColor(0xC9, 0x8A, 0x1C)

prs = Presentation()
prs.slide_width  = Inches(11)
prs.slide_height = Inches(8.5)
BLANK = prs.slide_layouts[6]

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    return slide

def add_text(slide, text, left, top, width, height,
             font_size=12, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb

def add_footer(slide, slide_num, total_slides=22, label=""):
    add_text(slide, f"{label}", 0.5, 8.10, 7.0, 0.30, font_size=9, color=DIM, italic=True)
    add_text(slide, f"{slide_num} / {total_slides}", 9.5, 8.10, 1.0, 0.30, font_size=9, color=DIM, align=PP_ALIGN.RIGHT)

def add_title_strip(slide, title, subtitle=None):
    add_text(slide, title, 0.5, 0.35, 10.0, 0.70, font_size=24, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name="Calibri")
    if subtitle:
        add_text(slide, subtitle, 0.5, 1.10, 10.0, 0.40, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)

def add_image(slide, path, left, top, width, height=None):
    if height is None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))

# ====================== SLIDE 1 — Title ======================
s = add_slide()
add_text(s, "Compositional monitoring of energy-mix drift on the simplex",
         0.5, 1.8, 10.0, 1.5, font_size=32, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "Five viewpoints. One observable stack. The hidden drivers, named.",
         0.5, 3.4, 10.0, 0.6, font_size=15, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "CoDaWork 2026  ·  11th International Workshop on Compositional Data Analysis",
         0.5, 4.7, 10.0, 0.4, font_size=13, color=DIM, align=PP_ALIGN.CENTER)
add_text(s, "Coimbra, Portugal  ·  1–5 June 2026",
         0.5, 5.1, 10.0, 0.4, font_size=12, color=DIM, align=PP_ALIGN.CENTER)
add_text(s, "P. Higgins  ·  Independent researcher  ·  Rogue Wave Audio, Markham, Ontario",
         0.5, 6.3, 10.0, 0.4, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "Manuscript available — see Standard Stamp on final slide.",
         0.5, 6.9, 10.0, 0.35, font_size=10, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_text(s, "The instrument reads.   The expert decides.   The hashes carry the receipts.   The vocabulary holds the line.",
         0.5, 7.7, 10.0, 0.35, font_size=10, italic=True, color=DIM, align=PP_ALIGN.CENTER)

# ====================== SLIDE 2 — The Question ======================
s = add_slide()
add_title_strip(s, "The Question",
                "What carriers really drove each country's energy transition — and which were doing the structural work?")
add_text(s, "When an electricity mix shifts over twenty-five years, two questions appear together:",
         0.7, 1.9, 9.6, 0.45, font_size=14, color=INK)
add_text(s, "1.   WHAT path did the system take?",
         1.0, 2.55, 9.0, 0.50, font_size=18, bold=True, color=GOLD)
add_text(s, "        Which carriers grew, which shrank, which stayed flat?",
         1.0, 3.05, 9.0, 0.40, font_size=13, color=DIM)
add_text(s, "2.   WHY did it take that path?",
         1.0, 3.7, 9.0, 0.50, font_size=18, bold=True, color=GOLD)
add_text(s, "        Which carriers were doing the structural work behind the change?",
         1.0, 4.20, 9.0, 0.40, font_size=13, color=DIM)
add_text(s,
         "The standard answer is a stacked-area chart of share-over-time.\nIt answers the first question and silently misleads the second.",
         0.7, 5.10, 9.6, 0.9, font_size=14, color=INK, italic=True)
add_text(s,
         "A small-share carrier whose share triples appears as a sliver.\nA large-share carrier drifting one point appears as the headline.",
         0.7, 6.20, 9.6, 0.9, font_size=13, color=DIM)
add_text(s,
         "This talk introduces the Activation Coefficient — the formal name for the carrier-level diagnostic that names the hidden drivers the size view never shows.",
         0.7, 7.30, 9.6, 0.5, font_size=12, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
add_footer(s, 2)

# ====================== SLIDE 3 — Size view hides the driver ======================
s = add_slide()
add_title_strip(s, "Viewpoint 1 — the size view",
                "What everyone sees when they look at the size view.  And what they miss.")
add_text(s, "WORLD ELECTRICITY GENERATION — the standard stacked-area chart",
         0.7, 1.85, 9.6, 0.45, font_size=14, bold=True, color=INK)
add_text(s,
         "•   Coal stays dominant for most of the period.\n"
         "•   Solar appears as a thin yellow sliver after 2010.\n"
         "•   Gas grows substantially.\n"
         "•   Nuclear declines slowly.\n"
         "•   Wind grows steadily but is small until late.",
         0.7, 2.4, 5.5, 2.0, font_size=13, color=INK)
add_text(s, "What the chart hides:",
         6.5, 1.95, 4.0, 0.45, font_size=15, bold=True, color=GOLD)
add_text(s,
         "USA, 2012–2013.\n\n"
         "Solar holds 0.107 % of the mix.\n"
         "Solar did 81.7 % of the structural work.\n\n"
         "Activation Coefficient:  760 × its size.",
         6.5, 2.5, 4.2, 3.0, font_size=14, color=INK)
add_text(s, "The size view did not show that.  No size view ever could.",
         6.5, 5.4, 4.2, 0.5, font_size=12, italic=True, color=ACCENT)
add_text(s,
         "The protocol presented today reads compositions natively in Aitchison geometry, stacks five viewpoints into one observable, and names the hidden drivers explicitly.",
         0.7, 6.4, 9.6, 1.0, font_size=13, color=INK)
add_text(s,
         "Mathematical foundation: standard compositional data analysis.   Operational claim: monitoring application is what is new.",
         0.7, 7.5, 9.6, 0.4, font_size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 3)

# ====================== SLIDE 4 — Five Viewpoints overview ======================
s = add_slide()
add_title_strip(s, "The five-viewpoint protocol",
                "Each viewpoint answers one specific question.  Combined, they yield a complete answer for any year-to-year transition step.")
add_image(s, MANU_FIG / "fig1_method.png", 1.7, 1.6, 7.6)
add_footer(s, 4)

# ====================== SLIDE 5 — Viewpoint 2: Helmsman ======================
s = add_slide()
add_title_strip(s, "Viewpoint 2 — the helmsman",
                "Who is at the wheel each year?  The single carrier with the largest CLR displacement at each transition step.")
add_text(s, "FORMAL DEFINITION  (Eq. 5)", 0.7, 1.9, 9.6, 0.45, font_size=13, bold=True, color=GOLD)
add_text(s, "σ(t)  =  argmax_i   | clr_i(t+1)  −  clr_i(t) |",
         0.7, 2.45, 9.6, 0.45, font_size=15, color=INK, italic=True, font_name="Cambria")
add_text(s,
         "The helmsman is a categorical assignment per step.  It names which carrier did the most directional work between two consecutive years.",
         0.7, 3.20, 9.6, 0.8, font_size=13, color=INK)
add_text(s, "READING THE HELMSMAN", 0.7, 4.40, 9.6, 0.4, font_size=13, bold=True, color=GOLD)
add_text(s,
         "•   The helmsman trajectory σ(1), σ(2), … shows who held the wheel year-by-year.\n"
         "•   The 'flip count' is the number of times the helmsman changes carrier across the window.\n"
         "•   Plotted with DOTTED lines — no continuous carrier-to-carrier path is implied.\n"
         "•   The helmsman is a categorical label, not a position on a continuum.",
         0.7, 4.90, 9.6, 2.4, font_size=13, color=INK)
add_text(s,
         "In our corpus the helmsman flip count ranges from 4 (World aggregate, stable) to 17 (Japan, post-Fukushima cascade).",
         0.7, 7.35, 9.6, 0.5, font_size=12, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_footer(s, 5)

# ====================== SLIDE 6 — Viewpoint 3: Power Share ======================
s = add_slide()
add_title_strip(s, "Viewpoint 3 — Power Share",
                "How much each carrier did at each transition.  Sums to 100 % across the carriers at every step.")
add_text(s, "FORMAL DEFINITION  (Eq. 6)", 0.7, 1.9, 9.6, 0.45, font_size=13, bold=True, color=GOLD)
add_text(s, "π_i(t)   =   ( Δ clr_i )²   /   Σ_j ( Δ clr_j )² ,        Σ_i  π_i(t)  =  1",
         0.7, 2.45, 9.6, 0.45, font_size=14, italic=True, color=INK, font_name="Cambria")
add_text(s,
         "The natural decomposition of the squared Aitchison distance across the carriers.  No carrier is hidden by the bookkeeping.",
         0.7, 3.15, 9.6, 0.8, font_size=13, color=INK)
add_text(s, "WHY IT MATTERS", 0.7, 4.30, 9.6, 0.4, font_size=13, bold=True, color=GOLD)
add_text(s,
         "•   The standard size view tells you what's BIG.\n"
         "•   The Power Share tells you what's MOVING.\n"
         "•   These are different questions.  A 30 %-share carrier holding steady has 0 % Power Share.\n"
         "•   A 0.1 %-share carrier in rapid growth can have 70–85 % Power Share at the step.",
         0.7, 4.80, 9.6, 2.4, font_size=13, color=INK)
add_text(s,
         "What's big ≠ what's moving.  This is the operational distinction that makes the rest of the protocol useful.",
         0.7, 7.35, 9.6, 0.5, font_size=12, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_footer(s, 6)

# ====================== SLIDE 7 — Viewpoint 4: Activation Coefficient ======================
s = add_slide()
add_title_strip(s, "Viewpoint 4 — the Activation Coefficient (the yeast factor)",
                "Why a small carrier mattered.  The leverage-to-size ratio — directional work divided by compositional weight.")
add_text(s, "FORMAL DEFINITION  (Eq. 7)", 0.7, 1.85, 9.6, 0.45, font_size=13, bold=True, color=GOLD)
add_text(s, "α_i(t)   =   π_i(t)   /   ρ_i(t)        reported when  ρ_i ≥ 0.1 %",
         0.7, 2.40, 9.6, 0.45, font_size=14, italic=True, color=INK, font_name="Cambria")
add_text(s, "Power Share divided by the carrier's composition share at the start of the step.",
         0.7, 3.05, 9.6, 0.4, font_size=12, color=DIM, italic=True)
add_text(s, "INTERPRETATION", 0.7, 3.65, 9.6, 0.4, font_size=13, bold=True, color=GOLD)
add_text(s,
         "α  =  1            the carrier did exactly its size's share of work\n"
         "α  ≫  1           a small carrier doing structural work far beyond its size  →  hidden driver  →  yeast\n"
         "α  ≪  1           a large carrier moving less than its size would predict   →  structural ballast",
         0.7, 4.10, 9.6, 1.8, font_size=13, color=INK, font_name="Calibri")
add_text(s, "WORKED EXAMPLE — USA Solar 2012 → 2013",
         0.7, 6.05, 9.6, 0.4, font_size=13, bold=True, color=GOLD)
add_text(s,
         "        composition share start of 2012  :   0.107 %        (solar was a tiny sliver)\n"
         "        power share of squared CLR motion :  81.7 %         (solar did most of the work)\n"
         "        activation coefficient           :   760 ×          (solar punched at 760× its size)",
         0.7, 6.55, 9.6, 1.2, font_size=12, color=INK, font_name="Consolas")
add_footer(s, 7)

# ====================== SLIDE 8 — Solar Yeast Era headline ======================
s = add_slide()
add_title_strip(s, "The hidden driver — solar, 2010–2015",
                "Across nine national electricity mixes, solar at sub-0.2 % composition share did 70–85 % of the structural directional work.")
add_text(s, "TOP YEAST MOMENTS  ·  9-COUNTRY CORPUS  ·  filtered to composition share ≥ 0.1 %",
         0.5, 1.85, 10.0, 0.4, font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
header_cols = ["#", "Country", "Transition", "Carrier", "AC × size", "Power Share", "Size at start"]
col_x = [0.5, 1.0, 1.85, 3.10, 5.30, 6.85, 8.55]
col_w = [0.4, 0.9, 1.3,  2.20, 1.55, 1.50, 1.45]
for i, h in enumerate(header_cols):
    add_text(s, h, col_x[i], 2.35, col_w[i], 0.40, font_size=11, bold=True, color=GOLD)

top_cases = [
    ("1", "USA", "2012–2013", "Solar",            "760 ×", "81.7 %", "0.107 %"),
    ("2", "FRA", "2010–2011", "Solar",            "659 ×", "72.6 %", "0.110 %"),
    ("3", "FRA", "2004–2005", "Wind",             "634 ×", "66.0 %", "0.104 %"),
    ("4", "CHN", "2013–2014", "Solar",            "549 ×", "84.7 %", "0.154 %"),
    ("5", "WLD", "2010–2011", "Solar",            "513 ×", "77.7 %", "0.151 %"),
    ("6", "FRA", "2005–2006", "Wind",             "495 ×", "83.5 %", "0.169 %"),
    ("7", "USA", "2013–2014", "Solar",            "395 ×", "87.9 %", "0.223 %"),
    ("8", "CHN", "2004–2005", "Bioenergy",        "346 ×", "39.7 %", "0.115 %"),
    ("9", "DEU", "2005–2006", "Solar",            "333 ×", "71.1 %", "0.214 %"),
    ("10","GBR", "2001–2002", "Wind",             "190 ×", "47.5 %", "0.249 %"),
]
y_base = 2.85
row_h = 0.42
for r, row in enumerate(top_cases):
    y_pos = y_base + r * row_h
    for c, val in enumerate(row):
        is_bold = (c == 1) or (c == 4)
        col = GOLD if c == 4 else INK
        add_text(s, val, col_x[c], y_pos, col_w[c], row_h, font_size=11, color=col, bold=is_bold, font_name="Consolas")

add_text(s, "Solar dominates seven of the top ten.  All in the 2010–2015 window.  Composition share between 0.1 % and 0.2 %.",
         0.5, 7.25, 10.0, 0.45, font_size=12, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "The size view never showed it.  The Activation Coefficient names it.",
         0.5, 7.75, 10.0, 0.35, font_size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
add_footer(s, 8)

# ====================== SLIDE 9 — Germany ======================
s = add_slide()
add_title_strip(s, "Germany — the continuous arc to the renewable vertex",
                "Twenty-five years of deliberate composition change.  The acceleration moments hidden in the size view.")
add_image(s, MANU_FIG / "fig2_germany.png", 1.5, 1.45, 8.0)
add_text(s,
         "Solar 2005–2006 at 0.21 % share, doing 71.1 % of the work — AC ≈ 333×.    The structural beginning of the Energiewende.",
         0.5, 7.85, 10.0, 0.30, font_size=10.5, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_footer(s, 9)

# ====================== SLIDE 10 — Japan ======================
s = add_slide()
add_title_strip(s, "Japan — Fukushima 2011 and the multi-year reorganisation",
                "An external shock that registers in every viewpoint.  The post-shock cascade tells the deeper story.")
add_image(s, MANU_FIG / "fig3_japan.png", 1.5, 1.45, 8.0)
add_text(s,
         "Helmsman flips: 17 — the loudest in the corpus.    Aitchison distance 2011 → 2012 = 3× neighbouring years.    Multi-year cascade.",
         0.5, 7.85, 10.0, 0.30, font_size=10.5, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_footer(s, 10)

# ====================== SLIDE 11 — United Kingdom ======================
s = add_slide()
add_title_strip(s, "United Kingdom — the coal exit as a regime change",
                "Policy-driven transition.  Coal goes from > 30 % to < 2 %.  Specific small carriers absorbed the displaced share.")
add_image(s, MANU_FIG / "fig4_uk.png", 1.5, 1.45, 8.0)
add_text(s,
         "Coal disappeared into specific, small, hard-working sub-categories — not into 'renewables' broadly.",
         0.5, 7.85, 10.0, 0.30, font_size=10.5, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_footer(s, 11)

# ====================== SLIDE 12 — Germany navigation chart ======================
s = add_slide()
add_title_strip(s, "Germany — the navigation chart",
                "PCA 2D projection of the CLR trajectory.  Continuous-arc archetype — course directness 0.41.")
add_image(s, MANU_FIG / "fig6_nav_deu.png", 0.7, 1.50, 9.6)
add_text(s,
         "Directional arc.  Long, deliberate motion through Aitchison space — the geometric signature of the Energiewende.",
         0.5, 7.85, 10.0, 0.30, font_size=10.5, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_footer(s, 12)

# ====================== SLIDE 13 — Japan navigation chart ======================
s = add_slide()
add_title_strip(s, "Japan — the navigation chart",
                "PCA 2D projection of the CLR trajectory.  Heavy-looping archetype — course directness 0.09.")
add_image(s, MANU_FIG / "fig6_nav_jpn.png", 0.7, 1.50, 9.6)
add_text(s,
         "Looping trajectory.  External shock (Fukushima 2011) and multi-year reorganisation register as path curvature.",
         0.5, 7.85, 10.0, 0.30, font_size=10.5, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_footer(s, 13)

# ====================== SLIDE 14 — United Kingdom navigation chart ======================
s = add_slide()
add_title_strip(s, "United Kingdom — the navigation chart",
                "PCA 2D projection of the CLR trajectory.  Jump-and-return archetype — course directness 0.36.")
add_image(s, MANU_FIG / "fig6_nav_gbr.png", 0.7, 1.50, 9.6)
add_text(s,
         "Coal-exit regime change.  The geometry records the displacement and where the share went.",
         0.5, 7.85, 10.0, 0.30, font_size=10.5, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_footer(s, 14)

# ====================== SLIDE 15 — Cross-country signature ======================
s = add_slide()
add_title_strip(s, "Cross-country signature — five of nine countries reproduce the deceptive-drift pattern",
                "Australia, China, United Kingdom, India, Japan fire the signature.  Germany, France, USA, World aggregate do not.")
add_image(s, MANU_FIG / "fig5_crosscountry.png", 1.0, 1.55, 9.0)
add_footer(s, 15)

# ====================== SLIDE 16 — Synthesis ======================
s = add_slide()
add_title_strip(s, "The synthesis  —  WHAT path  +  WHY",
                "Five viewpoints stack into one observable.  Each gives a partial glimpse; together they form the complete picture.")
row_defs = [
    ("Viewpoint 1  ·  Size view",             "WHAT carriers are big",       "Standard stacked-area chart"),
    ("Viewpoint 2  ·  Helmsman",              "WHO is at the wheel",         "Categorical per step"),
    ("Viewpoint 3  ·  Helmsman trajectory",   "WHEN the wheel changes",      "Dotted-line plot over time"),
    ("Viewpoint 4  ·  Power Share",           "HOW MUCH each carrier did",   "Squared-CLR decomposition (sums to 100 %)"),
    ("Viewpoint 5  ·  Activation Coefficient","WHY a small carrier mattered","Power Share ÷ composition share"),
]
y_start = 1.8
row_h = 0.95
for r, (vp, q, body) in enumerate(row_defs):
    y = y_start + r * row_h
    add_text(s, vp,   0.7, y, 3.4, 0.42, font_size=13, bold=True, color=GOLD)
    add_text(s, q,    4.2, y, 3.0, 0.42, font_size=13, bold=True, color=INK)
    add_text(s, body, 7.3, y, 3.2, 0.55, font_size=11.5, color=DIM)
add_text(s,
         "Used together, the five viewpoints answer  WHAT  ·  WHO  ·  WHEN  ·  HOW MUCH  ·  WHY  —  the complete answer.",
         0.5, 7.7, 10.0, 0.5, font_size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
add_footer(s, 16)

# ====================== SLIDE 17 — MC-4 falsifiable claim ======================
s = add_slide()
add_title_strip(s, "The falsifiable claim — MC-4 in three conjuncts",
                "Aitchison-native  +  formal change detection  +  carrier-level attribution  →  one observable stack.")
add_text(s, "Four explicit ways a CoDa specialist could defeat the claim:",
         0.7, 1.95, 9.6, 0.45, font_size=14, color=INK, italic=True)
def_paths = [
    ("1.  Prior-art defeat",  "Show the three-conjunct combination already exists in the literature.  Closest adjacent prior art identified: Morais, Thomas-Agnan & Simioni (2018); Arata & Onozaki (2017).  Neither combines all three to our reading; a closer match is welcome."),
    ("2.  Metric defeat",      "Show the verdicts reverse under a different valid simplex metric.  We have verified pair-invariance for TV distance and Aitchison distance across 101 datasets — but not exhausted the family."),
    ("3.  Case defeat",        "Show the 5-of-9 deceptive-drift signature is an artefact of preprocessing, carrier definition, or null-model choice rather than a robust compositional read.  Full corpus is public."),
    ("4.  Category defeat",    "Show compositional monitoring is most accurately an application note inside existing CoDa rather than a distinct monitoring category.  We hold no preconceived answer on this one."),
]
y_start = 2.55
row_h = 1.20
for r, (label, body) in enumerate(def_paths):
    y = y_start + r * row_h
    add_text(s, label, 0.7, y, 2.7, 0.40, font_size=13, bold=True, color=GOLD)
    add_text(s, body,  3.5, y, 7.0, row_h - 0.05, font_size=11, color=INK)
add_text(s,
         "The mathematics is not new.  The monitoring application may be.  If that sentence is wrong, this is the right room to kill it.",
         0.5, 7.65, 10.0, 0.5, font_size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_footer(s, 17)

# ====================== SLIDE 18 — Bridge to data scroll ======================
s = add_slide()
add_title_strip(s, "Now — every plate the engine produced",
                "What follows is not slides about data.  It is the data, run through the engine, scrolled through as a movie.")
add_text(s, "Nine countries.  Twenty-six years.  Six plates per country.",
         0.5, 2.4, 10.0, 0.5, font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "Master cover  →  9 country sections  ·  each section:",
         0.5, 3.3, 10.0, 0.5, font_size=13, color=INK, align=PP_ALIGN.CENTER)
add_text(s,
         "  ·  cover slide\n"
         "  ·  Stage 1 Section plate  —  bearings, CLR, plan view\n"
         "  ·  System course plot  —  CLR trajectory PCA (the navigation chart)\n"
         "  ·  Helmsman family  —  Sigma sequence, flip count, stability\n"
         "  ·  ILR-Helmert Triplet plate  —  three orthogonal scatter panels\n"
         "  ·  CNQ dashboard  —  quaternion-view diagnostics",
         2.0, 3.95, 7.5, 2.5, font_size=13, color=INK, font_name="Consolas")
add_text(s, "Sixty-six slides.  Pause me anywhere.  Every quantity hash-chained to the input CSV.",
         0.5, 6.9, 10.0, 0.5, font_size=14, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "→  CodaWork2026_PremierDataOutput_2026-05-13.pptx",
         0.5, 7.5, 10.0, 0.4, font_size=10.5, italic=True, color=DIM, align=PP_ALIGN.CENTER, font_name="Consolas")
add_footer(s, 18)

# ====================== SLIDE 19 — Bridge to projector + Q&A ======================
s = add_slide()
add_title_strip(s, "Q & A — with the manifold projector running",
                "Live HTML projector.  Interactive 3-D view of the energy-mix manifold.  Stays on through the discussion.")
add_text(s, "WHILE YOU ASK QUESTIONS",
         0.5, 2.3, 10.0, 0.5, font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "  ·  ORBIT  —  3-D rotation of the manifold\n"
         "  ·  TRAILS  —  show the trajectory paths over time\n"
         "  ·  LABELS  —  year labels along each path\n"
         "  ·  GHOST  —  partial-transparency overlay across countries\n"
         "  ·  FRONT / SIDE / TOP  —  fixed projections",
         2.0, 3.05, 7.5, 2.2, font_size=14, color=INK, font_name="Consolas")
add_text(s, "Run all  ·  Run AUS / CHN / DEU / FRA / GBR / IND / JPN / USA / WLD individually.",
         0.5, 5.4, 10.0, 0.5, font_size=12, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "Thank you.  I'll take questions.",
         0.5, 6.7, 10.0, 0.6, font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "→  codawork2026_projector.html   (open in browser; runs offline; no network required)",
         0.5, 7.55, 10.0, 0.4, font_size=10.5, italic=True, color=DIM, align=PP_ALIGN.CENTER, font_name="Consolas")
add_footer(s, 19)

# ====================== SLIDE 20 — Repositories + closing ======================
s = add_slide()
add_title_strip(s, "Repositories  ·  reproduce every plate in five minutes",
                "Hashes carry the receipts.  Anyone with the raw EMBER CSV can verify any plate.")
add_text(s, "higgins-decomposition (Hˢ)", 0.7, 2.1, 4.8, 0.5, font_size=15, bold=True, color=GOLD)
add_text(s,
         "Deterministic engines  ·  CNT v3.1.0  /  CNQ v2.0.0\n"
         "63 Investigation Catalog entries\n"
         "Hˢ Change Control v1.0  ·  HUF-STD-001/002/003\n"
         "Manuscript:  papers/codawork2026/manuscript/",
         0.7, 2.6, 5.0, 2.0, font_size=12, color=INK)
add_text(s, "github.com/PeterHiggins19/higgins-decomposition",
         0.7, 4.6, 5.0, 0.4, font_size=11, italic=True, color=GOLD, font_name="Consolas")

add_text(s, "Higgins-Unity-Framework (HUF)", 5.9, 2.1, 4.5, 0.5, font_size=15, bold=True, color=GOLD)
add_text(s,
         "MC-4 framing  ·  EITT canonical\n"
         "HUF Governance Charter\n"
         "KILL-001  ·  19 named failure modes\n"
         "Companion: case studies + cross-AI archive",
         5.9, 2.6, 4.5, 2.0, font_size=12, color=INK)
add_text(s, "github.com/PeterHiggins19/Higgins-Unity-Framework",
         5.9, 4.6, 5.0, 0.4, font_size=11, italic=True, color=GOLD, font_name="Consolas")

add_text(s, "Self-discipline note:  the L2 → TV correction is on the record.  We tell on ourselves when we mislabel.",
         0.7, 5.6, 9.6, 0.5, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_text(s, "The mathematics is not new.   The monitoring application may be.",
         0.7, 6.5, 9.6, 0.6, font_size=15, italic=True, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
add_text(s, "The talk is an ascent waypoint, not the summit.",
         0.7, 7.3, 9.6, 0.5, font_size=12, italic=True, color=INK, align=PP_ALIGN.CENTER)
add_footer(s, 20)

# ====================== SLIDE 21 — AI Use Declaration ======================
s = add_slide()
add_title_strip(s, "AI Use Declaration",
                "Per HUF Publication Standards (HUF-STD-001 v1.1) — conforming to scientific community conventions.")
add_text(s, "STANDARDS CONFORMANCE", 0.7, 1.95, 9.6, 0.4, font_size=12, bold=True, color=GOLD)
add_text(s,
         "ICMJE  ·  COPE  ·  Nature / Springer  ·  Science / AAAS  ·  WAME  ·  EU AI Act (2024)  ·  arXiv  ·  ACM  ·  IEEE",
         0.7, 2.40, 9.6, 0.5, font_size=12, color=INK)
add_text(s, "AI TOOLS USED", 0.7, 3.1, 9.6, 0.4, font_size=12, bold=True, color=GOLD)
add_text(s,
         "Claude (Anthropic)  ·  ChatGPT (OpenAI)  ·  Copilot (Microsoft)  ·  Gemini (Google)  ·  Grok (xAI)",
         0.7, 3.55, 9.6, 0.5, font_size=12, color=INK)
add_text(s, "Collectively: the HUF AI Collective.", 0.7, 4.0, 9.6, 0.4, font_size=11, italic=True, color=DIM)
add_text(s, "TASKS PERFORMED BY AI", 0.7, 4.6, 9.6, 0.4, font_size=12, bold=True, color=GOLD)
add_text(s,
         "Drafting of supporting documents and slides  ·  consistency editing  ·  cross-checking claims across documents\n"
         "Literature-search assistance for prior-art areas  ·  adversarial review for overclaim and drift catching",
         0.7, 5.05, 9.6, 0.9, font_size=11, color=INK)
add_text(s, "AUTHOR RESPONSIBILITY", 0.7, 6.05, 9.6, 0.4, font_size=12, bold=True, color=GOLD)
add_text(s,
         "P. Higgins retains full responsibility for all scientific claims, data interpretation, methodological choices,\n"
         "and conclusions.  All AI-generated content has been reviewed and verified.  AI tools are NOT listed as authors.",
         0.7, 6.50, 9.6, 1.0, font_size=11, color=INK)
add_footer(s, 21)

# ====================== SLIDE 22 — Standard Stamp colophon ======================
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = RGBColor(0x14, 0x2B, 0x46)
add_text(s, "Hˢ Framework  ·  Engine Declaration",
         0.5, 0.35, 10.0, 0.6, font_size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "Standard Stamp — appended to every Hˢ-produced document  ·  HUF-STD-001 compliant",
         0.5, 0.95, 10.0, 0.35, font_size=11, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "From  «  CoDaWork 2026 — Final Talk — 2026-05-17  »",
         0.5, 1.40, 10.0, 0.35, font_size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)

add_text(s, "The Framework", 0.6, 1.95, 3.2, 0.4, font_size=13, bold=True, color=GOLD)
add_text(s,
         "Higgins Decomposition (Hˢ)\n"
         "  Deterministic compositional inference\n"
         "  Operates on the simplex in Aitchison geometry\n"
         "\n"
         "Principle\n"
         "  CNT measures invariance.\n"
         "  CNQ names the algebra it lives in.\n"
         "\n"
         "Open  ·  Hash-chained  ·  Reproducible\n"
         "  Apache-2.0 code  +  CC BY 4.0 docs\n"
         "  Same input → same output, always.\n"
         "\n"
         "Validated across\n"
         "  11 domains  ·  101 reference datasets\n"
         "  44 orders of magnitude in scale\n"
         "  Three IEEE-floor confirmations:\n"
         "  drive failures, CMB photons,\n"
         "  Standard-Model neutrino oscillation.",
         0.6, 2.40, 3.4, 5.5, font_size=9.5, color=INK)

add_text(s, "Engines  ·  Methods", 3.95, 1.95, 3.4, 0.4, font_size=13, bold=True, color=GOLD)
add_text(s,
         "CNT  ·  Compositional Navigation Tensor\n"
         "  v3.1.0  (schema 3.1.0)\n"
         "  cnt.py  (Python)  +  cnt.R  (R port)\n"
         "  Closure → CLR → ILR-Helmert →\n"
         "  bearing tensor, angular velocity,\n"
         "  Higgins scale, depth tower,\n"
         "  IR class, attractor fit.\n"
         "\n"
         "CNQ  ·  Compositional Navigation Quaternion\n"
         "  v2.0.0  (schema cnq/2.0.0)\n"
         "  cnq.py (Python)  +  cnq.R (R port)\n"
         "  Quaternion view  ·  CHSH joint-coherence\n"
         "  Twin-quaternion factoring at D=8\n"
         "\n"
         "Plate suite  ·  HUF-STD-002 link 4\n"
         "  Stage 0  Foundations  (HUF-STD-003)\n"
         "  Stage 1  Section + ILR-Helmert Triplet\n"
         "  Stage 2  Helmsman + course + CoDa-PCA\n"
         "  Stage 3  Depth tower + attractor + κ^Hs\n"
         "  Outputs:  PDF  /  PNG  /  SVG",
         3.95, 2.40, 3.4, 5.5, font_size=9.5, color=INK)

add_text(s, "Find us  ·  Contact", 7.4, 1.95, 3.3, 0.4, font_size=13, bold=True, color=GOLD)
add_text(s,
         "Repository\n"
         "  github.com/PeterHiggins19/\n"
         "  higgins-decomposition\n"
         "\n"
         "Quick start (AI or researcher)\n"
         "  Load HS_FAST_REFRESH.json first\n"
         "  Read AI_AGENTS.md for the protocol\n"
         "  Reproduce with one command\n"
         "  per REPRODUCIBILITY_CHECKLIST.md\n"
         "\n"
         "Sibling repositories\n"
         "  Higgins-Unity-Framework (HUF parent)\n"
         "  Rogue-Wave-Audio (DADC lineage)\n"
         "\n"
         "Author  ·  Lab\n"
         "  Peter Higgins\n"
         "  Rogue Wave Audio /\n"
         "  Binaural Test Lab\n"
         "  Markham, Ontario, Canada\n"
         "  RWA-001  (lab identity card)\n"
         "\n"
         "Contact\n"
         "  PeterHiggins@RogueWaveAudio.com\n"
         "  Open an issue on GitHub\n"
         "  Help is available — free, no gatekeeping.",
         7.4, 2.40, 3.3, 5.5, font_size=9.5, color=INK)

add_text(s, "The instrument reads.   The expert decides.   The hashes carry the receipts.   The vocabulary holds the line.",
         0.5, 7.75, 10.0, 0.35, font_size=10, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "Conforms to  HUF-STD-001 (Publication Standards)  ·  HUF-STD-002 (Tensor Train I/O)  ·  HUF-STD-003 (Linear Algebra Foundations)",
         0.5, 8.10, 10.0, 0.30, font_size=9, color=DIM, align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"OK: {OUT}")
print(f"slides: {len(prs.slides)}")
