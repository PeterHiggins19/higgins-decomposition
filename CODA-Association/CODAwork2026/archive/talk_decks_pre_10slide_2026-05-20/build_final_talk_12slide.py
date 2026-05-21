"""
CoDaWork 2026 — Final Talk Deck builder, COMPRESSED 12-SLIDE VERSION.

Built 2026-05-20 from ChatGPT-prepared compression plan
(`CodaWork2026_FinalTalk_12Slide_CompressionPlan.json`).
DOES NOT supersede the 22-slide deck — it ships as a sibling.

Story arc (12 slides):
   1. Title + question
   2. Size view hides work        (USA Solar 760× hook)
   3. Five viewpoints              (one schematic)
   4. Activation Coefficient       (yeast factor + formula)
   5. Three archetypes overview    (DEU / JPN / GBR)
   6. Germany — continuous arc
   7. Japan — shock + reorganisation
   8. UK — regime change
   9. 5-of-9 cross-country signature
  10. What the stack answers       (synthesis)
  11. MC-4 falsifiable claim
  12. Inspect the instrument       (scroll + projector + repo + AI use footer)

Total live time target: 10.5 min talk + 1.5 min scroll + 1 min projector
                       + 2 min buffer = ~15 min.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUT = Path("/sessions/epic-gracious-lovelace/mnt/Claude CoWorker/Current-Repo/Hs/CODA-Association/CODAwork2026/data_outputs/CodaWork2026_FinalTalk_12Slide_2026-05-20.pptx")
MANU_FIG = Path("/sessions/epic-gracious-lovelace/mnt/Claude CoWorker/Current-Repo/Hs/papers/codawork2026/manuscript/figures")

# Palette — preserved from 22-slide deck for visual continuity
NAVY   = RGBColor(0x0B, 0x1F, 0x33)
GOLD   = RGBColor(0xF2, 0xB6, 0x32)
INK    = RGBColor(0xEE, 0xEE, 0xEE)
DIM    = RGBColor(0xB8, 0xB8, 0xB8)
ACCENT = RGBColor(0xC9, 0x8A, 0x1C)

prs = Presentation()
prs.slide_width  = Inches(11)
prs.slide_height = Inches(8.5)
BLANK = prs.slide_layouts[6]
TOTAL = 12

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    return slide

def add_text(slide, text, left, top, width, height,
             font_size=12, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb

def add_footer(slide, n, label=""):
    add_text(slide, label, 0.5, 8.10, 7.0, 0.30, font_size=9, color=DIM, italic=True)
    add_text(slide, f"{n} / {TOTAL}", 9.5, 8.10, 1.0, 0.30, font_size=9, color=DIM, align=PP_ALIGN.RIGHT)

def add_title_strip(slide, title, subtitle=None):
    add_text(slide, title, 0.5, 0.35, 10.0, 0.70, font_size=26, bold=True, color=INK, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, 0.5, 1.10, 10.0, 0.40, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)

def add_image(slide, path, left, top, width, height=None):
    if height is None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))

# ───────── SLIDE 1 — Title + question ─────────
s = add_slide()
add_text(s, "Compositional monitoring of energy-mix drift on the simplex",
         0.5, 1.7, 10.0, 1.5, font_size=30, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "Which carrier did the structural work?",
         0.5, 3.4, 10.0, 0.6, font_size=18, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Not only which carrier got bigger — which carrier moved the composition.",
         0.5, 4.2, 10.0, 0.5, font_size=13, color=INK, italic=True, align=PP_ALIGN.CENTER)
add_text(s, "CoDaWork 2026  ·  Coimbra, Portugal  ·  1–5 June 2026",
         0.5, 5.7, 10.0, 0.4, font_size=13, color=DIM, align=PP_ALIGN.CENTER)
add_text(s, "P. Higgins  ·  Rogue Wave Audio / Binaural Test Lab  ·  Markham, Ontario, Canada",
         0.5, 6.2, 10.0, 0.4, font_size=12, color=DIM, align=PP_ALIGN.CENTER)
add_text(s,
         "Operationalizing compositional analysis — a runnable standard for researchers and the AI assistants they choose.",
         0.5, 7.1, 10.0, 0.4, font_size=11, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "The instrument reads.   The expert decides.   The hashes carry the receipts.   The vocabulary holds the line.",
         0.5, 7.7, 10.0, 0.35, font_size=10, italic=True, color=DIM, align=PP_ALIGN.CENTER)

# ───────── SLIDE 2 — Size view hides work ─────────
s = add_slide()
add_title_strip(s, "The size view hides the work",
                "A carrier can be small in share and large in structural work.")
add_text(s, "World electricity, 25 years — the standard stacked-area view",
         0.7, 1.85, 9.6, 0.4, font_size=14, bold=True, color=INK)
add_text(s,
         "Coal stays dominant.  Gas grows.  Nuclear declines.\n"
         "Solar appears as a thin yellow sliver after 2010.\n"
         "Wind grows steadily but stays visually small until late.",
         0.7, 2.35, 5.4, 1.6, font_size=13, color=INK)
add_text(s, "What the size view misses",
         6.6, 1.95, 4.0, 0.45, font_size=15, bold=True, color=GOLD)
add_text(s,
         "USA Solar, 2012 → 2013",
         6.6, 2.55, 4.2, 0.45, font_size=14, bold=True, color=INK)
add_text(s,
         "starting share        0.107 %\n"
         "structural Power Share  81.7 %\n"
         "Activation Coefficient   ≈ 760 ×",
         6.6, 3.05, 4.2, 1.5, font_size=14, color=INK, font_name="Consolas")
add_text(s,
         "Solar acted at 760 × its size.\nNo size view shows that.",
         6.6, 4.85, 4.2, 1.0, font_size=13, italic=True, color=ACCENT)
add_text(s,
         "This talk is the reason that number exists.",
         0.7, 6.6, 9.6, 0.5, font_size=14, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
add_text(s,
         "Mathematics: standard compositional data analysis.  Application: monitoring frame may be new.",
         0.7, 7.5, 9.6, 0.4, font_size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 2)

# ───────── SLIDE 3 — Five viewpoints, one observable stack ─────────
s = add_slide()
add_title_strip(s, "Five viewpoints, one observable stack",
                "Each viewpoint answers one question.  Together: an auditable transition event.")
add_image(s, MANU_FIG / "fig1_method.png", 0.6, 1.7, 5.6)
add_text(s, "Composition",      6.5, 1.85, 4.0, 0.40, font_size=15, bold=True, color=GOLD)
add_text(s, "what share each carrier has.", 6.5, 2.25, 4.0, 0.35, font_size=12, color=INK)
add_text(s, "Helmsman",         6.5, 2.75, 4.0, 0.40, font_size=15, bold=True, color=GOLD)
add_text(s, "which carrier has the largest CLR move at a step.", 6.5, 3.15, 4.0, 0.4, font_size=12, color=INK)
add_text(s, "Helmsman trajectory", 6.5, 3.75, 4.0, 0.40, font_size=15, bold=True, color=GOLD)
add_text(s, "when the steering carrier changes.", 6.5, 4.15, 4.0, 0.35, font_size=12, color=INK)
add_text(s, "Power Share",      6.5, 4.65, 4.0, 0.40, font_size=15, bold=True, color=GOLD)
add_text(s, "how much squared CLR motion each carrier did.", 6.5, 5.05, 4.0, 0.4, font_size=12, color=INK)
add_text(s, "Activation Coefficient", 6.5, 5.65, 4.0, 0.40, font_size=15, bold=True, color=GOLD)
add_text(s, "Power Share ÷ starting share — the yeast factor.", 6.5, 6.05, 4.0, 0.4, font_size=12, color=INK)
add_text(s,
         "All five derive from CLR + ILR-Helmert.  Pure CoDa geometry; no new mathematics.",
         0.5, 7.5, 10.0, 0.4, font_size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 3)

# ───────── SLIDE 4 — Activation Coefficient, the yeast factor ─────────
s = add_slide()
add_title_strip(s, "The Activation Coefficient — the yeast factor",
                "How much structural work a carrier does, relative to how much of the mix it is.")
add_text(s, "α_i(t)  =  Power Share_i(t)  ÷  starting share_i(t)",
         0.5, 1.95, 10.0, 0.7, font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font_name="Consolas")
add_text(s, "α ≈ 1     carrier does work proportional to its size — ordinary",
         1.0, 3.0, 9.0, 0.45, font_size=14, color=INK)
add_text(s, "α ≫ 1     carrier acts far above its size — hidden driver",
         1.0, 3.5, 9.0, 0.45, font_size=14, color=GOLD, bold=True)
add_text(s, "α < 1     carrier carries less work than its size suggests — coasting",
         1.0, 4.0, 9.0, 0.45, font_size=14, color=DIM)
add_text(s, "Worked example — USA Solar 2012 → 2013",
         0.7, 4.85, 9.6, 0.45, font_size=15, bold=True, color=INK)
add_text(s,
         "starting share      0.107 %     small\n"
         "Power Share         81.7  %     most of the work\n"
         "α                   ≈ 760 ×     yeast",
         0.7, 5.35, 9.6, 1.5, font_size=14, color=INK, font_name="Consolas")
add_text(s,
         "Yeast is 2% of a loaf by mass and does 100% of the rising. Same shape.",
         0.7, 7.05, 9.6, 0.45, font_size=12, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s,
         "Solar 2010–2015 appears repeatedly as small-share / large-structural-work across the corpus.",
         0.7, 7.55, 9.6, 0.4, font_size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 4)

# ───────── SLIDE 5 — Three archetypes ─────────
s = add_slide()
add_title_strip(s, "Three archetypes — one instrument, three regimes",
                "Same protocol applied to three transitions that look fundamentally different.")
add_text(s, "Germany",   1.05, 1.85, 3.0, 0.5, font_size=18, bold=True, color=GOLD)
add_text(s, "deliberate transition",  1.05, 2.30, 3.0, 0.4, font_size=12, italic=True, color=INK)
add_text(s, "continuous arc",         1.05, 2.65, 3.0, 0.4, font_size=12, italic=True, color=DIM)
add_text(s, "Energiewende\n2000 → 2025\nsolar + wind absorb\nstructural work\nbefore size dominates.",
         1.05, 3.25, 3.0, 3.0, font_size=13, color=INK)

add_text(s, "Japan",     4.5, 1.85, 3.0, 0.5, font_size=18, bold=True, color=GOLD)
add_text(s, "external shock",         4.5, 2.30, 3.0, 0.4, font_size=12, italic=True, color=INK)
add_text(s, "loop and reorganise",    4.5, 2.65, 3.0, 0.4, font_size=12, italic=True, color=DIM)
add_text(s, "Fukushima 2011\ndisplaces nuclear,\ncauses 2011–2013\nmulti-year compositional\nreorganisation.",
         4.5, 3.25, 3.0, 3.0, font_size=13, color=INK)

add_text(s, "United Kingdom", 7.95, 1.85, 3.0, 0.5, font_size=18, bold=True, color=GOLD)
add_text(s, "regime change",          7.95, 2.30, 3.0, 0.4, font_size=12, italic=True, color=INK)
add_text(s, "jump and return",        7.95, 2.65, 3.0, 0.4, font_size=12, italic=True, color=DIM)
add_text(s, "Coal exit 2012–2020\nfrom > 30 % to < 2 %.\nWind, solar, others\nabsorb displaced\nstructural work.",
         7.95, 3.25, 3.0, 3.0, font_size=13, color=INK)

add_text(s,
         "Three different transition regimes.  One operational protocol reads them all.",
         0.5, 7.4, 10.0, 0.5, font_size=13, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
add_footer(s, 5)

# ───────── SLIDE 6 — Germany ─────────
s = add_slide()
add_title_strip(s, "Germany — deliberate transition as continuous course",
                "The Energiewende read as a single smooth arc on the simplex.")
add_image(s, MANU_FIG / "fig2_germany.png", 0.5, 1.55, 7.6, height=5.0)
add_image(s, MANU_FIG / "fig6_nav_deu.png", 8.25, 1.55, 2.6)
add_text(s,
         "Solar 2005–2006:  0.21 % share  ·  71.1 % structural work  ·  α ≈ 333 ×",
         0.5, 6.75, 10.0, 0.4, font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Helmsman trajectory: course directness 0.41 — continuous arc to renewable vertex.",
         0.5, 7.20, 10.0, 0.4, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 6)

# ───────── SLIDE 7 — Japan ─────────
s = add_slide()
add_title_strip(s, "Japan — Fukushima shock and reorganisation",
                "The instrument detects both the shock and the multi-year reorganisation.")
add_image(s, MANU_FIG / "fig3_japan.png", 0.5, 1.55, 7.6, height=5.0)
add_image(s, MANU_FIG / "fig6_nav_jpn.png", 8.25, 1.55, 2.6)
add_text(s,
         "Aitchison distance 2011 → 2012  ≈ 3 × neighbouring-year baseline  ·  helmsman flips 17 ×",
         0.5, 6.75, 10.0, 0.4, font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Course directness 0.09 — looping reorganisation, not a single step.",
         0.5, 7.20, 10.0, 0.4, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 7)

# ───────── SLIDE 8 — UK ─────────
s = add_slide()
add_title_strip(s, "United Kingdom — coal exit as regime change",
                "Policy-driven displacement, absorbed across multiple renewable carriers.")
add_image(s, MANU_FIG / "fig4_uk.png", 0.5, 1.55, 7.6, height=5.0)
add_image(s, MANU_FIG / "fig6_nav_gbr.png", 8.25, 1.55, 2.6)
add_text(s,
         "Coal:  > 30 %  →  < 2 %.   Wind, solar, and other renewables absorb the displaced structural work.",
         0.5, 6.75, 10.0, 0.4, font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Course directness 0.36 — jump-and-return archetype.",
         0.5, 7.20, 10.0, 0.4, font_size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER)
add_footer(s, 8)

# ───────── SLIDE 9 — Cross-country: 5 of 9 ─────────
s = add_slide()
add_title_strip(s, "Cross-country signature — 5 of 9 reproduce the pattern",
                "From three case archetypes to a corpus-level result.")
add_image(s, MANU_FIG / "fig5_crosscountry.png", 1.7, 1.55, 7.6, height=5.0)
add_text(s, "fires   AUS · CHN · GBR · IND · JPN",
         0.5, 7.10, 10.0, 0.45, font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "quiet   DEU (annual) · FRA · USA · WLD",
         0.5, 7.50, 10.0, 0.4, font_size=12, color=DIM, align=PP_ALIGN.CENTER)
add_text(s,
         "A useful detector should not fire everywhere.  Discrimination is itself evidence.",
         0.5, 7.95, 10.0, 0.4, font_size=11, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_footer(s, 9)

# ───────── SLIDE 10 — What the stack answers ─────────
s = add_slide()
add_title_strip(s, "What the stack answers",
                "One observable, five distinct questions, one reproducible object.")
rows = [
    ("WHAT", "carriers are big.",                    "size view"),
    ("WHO",  "is at the wheel.",                     "Helmsman"),
    ("WHEN", "the steering changes.",                "Helmsman trajectory"),
    ("HOW MUCH", "work each carrier did.",           "Power Share"),
    ("WHY",  "a small carrier mattered.",            "Activation Coefficient"),
]
y0 = 2.0
for i, (q, ans, who) in enumerate(rows):
    y = y0 + i * 0.85
    add_text(s, q,   0.9, y, 1.7, 0.6, font_size=20, bold=True, color=GOLD)
    add_text(s, ans, 2.7, y, 4.6, 0.6, font_size=15, color=INK)
    add_text(s, who, 7.3, y, 3.2, 0.6, font_size=13, italic=True, color=DIM)
add_text(s,
         "The stack does not replace interpretation.  It gives interpretation a reproducible object.",
         0.5, 7.5, 10.0, 0.5, font_size=13, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_footer(s, 10)

# ───────── SLIDE 11 — MC-4 falsifiable claim ─────────
s = add_slide()
add_title_strip(s, "MC-4 — the falsifiable claim",
                "Three conjuncts.  Four defeat paths.  One open invitation to the room.")
add_text(s, "The claim", 0.7, 1.85, 4.5, 0.45, font_size=16, bold=True, color=GOLD)
add_text(s,
         "Hˢ is the first published protocol that combines:\n\n"
         "  1.  Aitchison-native compositional metrics\n"
         "  2.  formal change-point / deceptive-drift detection\n"
         "  3.  carrier-level structural-work attribution",
         0.7, 2.35, 5.0, 3.5, font_size=13, color=INK)
add_text(s, "Defeat paths (any one falsifies)", 6.0, 1.85, 4.5, 0.45, font_size=16, bold=True, color=GOLD)
add_text(s,
         "  A.   prior-art path — show all three were already combined.\n\n"
         "  B.   metric path — show a defensible non-Aitchison metric.\n\n"
         "  C.   case path — show a country / year the protocol misreads.\n\n"
         "  D.   category path — show the construct itself is ill-posed.",
         6.0, 2.35, 4.7, 4.5, font_size=12, color=INK)
add_text(s,
         "This is the claim I am asking the CoDa community to test.  Take it apart.",
         0.5, 7.5, 10.0, 0.5, font_size=13, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
add_footer(s, 11, "INV-051 CANONICAL · 5 of 9 reproduces empirically")

# ───────── SLIDE 12 — Inspect the instrument ─────────
s = add_slide()
add_title_strip(s, "Now inspect the instrument",
                "I will not ask you to trust the speaker.  The outputs are open for inspection.")

add_text(s, "Manuscript",   1.0, 1.95, 3.0, 0.5, font_size=17, bold=True, color=GOLD)
add_text(s,
         "25-page paper,\nNature structure,\nSupplementary Info,\nsix figures.\n\n"
         "papers/codawork2026/\n  manuscript/",
         1.0, 2.50, 3.0, 3.5, font_size=12, color=INK)

add_text(s, "Cinema scroll", 4.3, 1.95, 3.0, 0.5, font_size=17, bold=True, color=GOLD)
add_text(s,
         "66-slide reel,\n325-page master PDF.\nEvery plate the engine\nproduced for the\nnine-country EMBER\ncorpus.\nScroll at speed.\nPause anywhere.",
         4.3, 2.50, 3.0, 3.5, font_size=12, color=INK)

add_text(s, "Projector",    7.6, 1.95, 3.0, 0.5, font_size=17, bold=True, color=GOLD)
add_text(s,
         "Offline HTML.\nRADAR · BARY · ALIGN\n+ SHOCK overlay.\n\nQ&A backdrop.\n\nNo network call.\nRuns in any browser.",
         7.6, 2.50, 3.0, 3.5, font_size=12, color=INK)

add_text(s,
         "github.com / PeterHiggins19 / higgins-decomposition   ·   CODA-Association/   ·   QR on handout",
         0.5, 6.40, 10.0, 0.4, font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s,
         "Hand-out in UN-6 locales (EN · FR · ES · RU · ZH · AR).  Operationalization standard, runnable.",
         0.5, 6.85, 10.0, 0.4, font_size=11, italic=True, color=INK, align=PP_ALIGN.CENTER)

add_text(s,
         "AI Use Declaration (HUF-STD-001 v1.1): research design, mathematical content, code, and scientific responsibility "
         "remain with the named author.  AI assistants (Claude, ChatGPT, Grok) used for drafting, sweeps, and reviews.  "
         "Author retains full responsibility.   Apache-2.0 code  ·  CC BY 4.0 docs.",
         0.5, 7.35, 10.0, 0.7, font_size=9, color=DIM, align=PP_ALIGN.CENTER)
add_text(s,
         "The instrument reads.  The expert decides.  The hashes carry the receipts.  "
         "The vocabulary holds the line.  The AI follows the same protocol.",
         0.5, 8.05, 10.0, 0.35, font_size=10, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# Save
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
